"""
Custom document reading tools for ProfileMatch
These tools directly read documents without requiring vector databases.
"""

from crewai.tools import BaseTool
from typing import Type, Optional
from pydantic import BaseModel, Field
import PyPDF2
import docx
from pptx import Presentation
from pathlib import Path


class DocumentInput(BaseModel):
    """Input schema for document reading tools."""
    file_path: str = Field(..., description="Full path to the document file to read")


class PDFReaderTool(BaseTool):
    name: str = "Read PDF Document"
    description: str = (
        "Reads and extracts all text content from a PDF file. "
        "Provide the full file path as input. "
        "Returns the complete text content of the PDF."
    )
    args_schema: Type[BaseModel] = DocumentInput

    def _run(self, file_path: str) -> str:
        """Read a PDF file and return its text content."""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            
            if not text.strip():
                return f"PDF file {file_path} appears to be empty or could not be read."
            
            return text.strip()
        except Exception as e:
            return f"Error reading PDF {file_path}: {str(e)}"


class DOCXReaderTool(BaseTool):
    name: str = "Read Word Document"
    description: str = (
        "Reads and extracts all text content from a Word (DOCX) file. "
        "Provide the full file path as input. "
        "Returns the complete text content of the document."
    )
    args_schema: Type[BaseModel] = DocumentInput

    def _run(self, file_path: str) -> str:
        """Read a DOCX file and return its text content."""
        try:
            doc = docx.Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            if not text.strip():
                return f"Word document {file_path} appears to be empty."
            
            return text.strip()
        except Exception as e:
            return f"Error reading Word document {file_path}: {str(e)}"


class PPTXReaderTool(BaseTool):
    name: str = "Read PowerPoint Presentation"
    description: str = (
        "Reads and extracts all text content from a PowerPoint (PPTX) file. "
        "Provide the full file path as input. "
        "Returns the complete text content from all slides."
    )
    args_schema: Type[BaseModel] = DocumentInput

    def _run(self, file_path: str) -> str:
        """Read a PPTX file and return its text content."""
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n--- Slide {slide_num} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            
            if not text.strip():
                return f"PowerPoint file {file_path} appears to be empty or has no text content."
            
            return text.strip()
        except Exception as e:
            return f"Error reading PowerPoint {file_path}: {str(e)}"
